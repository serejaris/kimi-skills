import os
from pathlib import Path

# Try to import advanced libraries, fallback to simple estimation if not available
try:
    import tiktoken
    TIKTOKEN_AVAILABLE = True
except ImportError:
    TIKTOKEN_AVAILABLE = False
    print("Warning: tiktoken not available, using character-based chunking")

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

# Multiple PDF parsers for robustness
try:
    from pypdf import PdfReader as PyPdfReader
    PYPDF_AVAILABLE = True
except ImportError:
    PYPDF_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    from pdfminer.high_level import extract_text as pdfminer_extract
    PDFMINER_AVAILABLE = True
except ImportError:
    PDFMINER_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class DocumentSplitter:
    def __init__(self, output_dir="/mnt/agents/chunks", model_name="cl100k_base", chunk_size=32000, overlap_ratio=0.1):
        """
        初始化切分工具
        :param output_dir: 切分后的文件保存目录
        :param model_name: 使用的分词模型 (默认 cl100k_base 适配 GPT-4/3.5)
        :param chunk_size: 每个 Chunk 的 Token 上限 (默认 64k)
        :param overlap_ratio: 块之间的重叠比例 (默认 0.1 即 10%)
        """
        self.output_dir = Path(output_dir)
        self.chunk_size = chunk_size
        self.overlap_ratio = overlap_ratio

        if TIKTOKEN_AVAILABLE:
            self.encoding = tiktoken.get_encoding(model_name)
            self.use_tiktoken = True
            self.overlap_size = int(chunk_size * overlap_ratio)
        else:
            self.encoding = None
            self.use_tiktoken = False
            # Fallback: 1 token ≈ 4 characters
            self.chunk_size_chars = chunk_size * 4
            self.overlap_size_chars = int(self.chunk_size_chars * overlap_ratio)

        # 确保输出目录存在
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def process_file(self, file_path):
        """
        处理单个文件：提取 -> 切分 -> 保存
        :return: 生成的文件绝对路径列表
        """
        file_path = Path(file_path)
        if not file_path.exists():
            print(f"Error: File not found {file_path}")
            return []

        # 1. 提取文本
        print(f"Reading file: {file_path.name} ...")
        full_text = self._extract_text(file_path)

        if not full_text:
            print(f"Warning: File content is empty or cannot be extracted: {file_path.name}")
            return []

        # 2. Tokenize or use character-based chunking
        generated_files = []
        base_name = file_path.stem  # 文件名（无后缀）

        if self.use_tiktoken:
            # Use tiktoken for accurate token-based chunking with overlap
            tokens = self.encoding.encode(full_text)
            total_tokens = len(tokens)

            print(f"  - Total tokens: {total_tokens}")
            print(f"  - Chunk size: {self.chunk_size} tokens, overlap: {self.overlap_size} tokens ({self.overlap_ratio*100:.0f}%)")

            part_num = 0
            i = 0
            while i < total_tokens:
                # Calculate chunk with overlap
                chunk_end = min(i + self.chunk_size, total_tokens)
                chunk_tokens = tokens[i : chunk_end]
                chunk_text = self.encoding.decode(chunk_tokens)

                part_num += 1
                output_filename = f"{base_name}_part_{part_num}.txt"
                output_path = self.output_dir / output_filename

                try:
                    with open(output_path, "w", encoding="utf-8") as f:
                        f.write(chunk_text)
                    generated_files.append(str(output_path.absolute()))
                except Exception as e:
                    print(f"  - Failed to write chunk {part_num}: {e}")

                # Move to next chunk with overlap (step = chunk_size - overlap)
                i += (self.chunk_size - self.overlap_size)

            print(f"  - Actual number of chunks: {part_num}")

        else:
            # Fallback: character-based chunking with overlap
            total_chars = len(full_text)

            print(f"  - Total characters: {total_chars} (estimated ~{total_chars // 4} tokens)")
            print(f"  - Chunk size: {self.chunk_size_chars} chars, overlap: {self.overlap_size_chars} chars ({self.overlap_ratio*100:.0f}%)")

            part_num = 0
            i = 0
            while i < total_chars:
                # Calculate chunk with overlap
                chunk_end = min(i + self.chunk_size_chars, total_chars)
                chunk_text = full_text[i : chunk_end]

                part_num += 1
                output_filename = f"{base_name}_part_{part_num}.txt"
                output_path = self.output_dir / output_filename

                try:
                    with open(output_path, "w", encoding="utf-8") as f:
                        f.write(chunk_text)
                    generated_files.append(str(output_path.absolute()))
                except Exception as e:
                    print(f"  - Failed to write chunk {part_num}: {e}")

                # Move to next chunk with overlap (step = chunk_size - overlap)
                i += (self.chunk_size_chars - self.overlap_size_chars)

            print(f"  - Actual number of chunks: {part_num}")

        return generated_files

    # --- 以下是复用的文本提取逻辑 ---
    
    def _extract_text(self, file_path):
        ext = file_path.suffix.lower()
        try:
            if ext == '.pdf':
                if PYMUPDF_AVAILABLE or PYPDF_AVAILABLE or PDFMINER_AVAILABLE:
                    return self._extract_pdf(file_path)
                else:
                    print(f"  - PDF extraction requires PyMuPDF, pypdf, or pdfminer (none installed)")
                    return ""
            elif ext in ['.docx', '.doc']:
                if DOCX_AVAILABLE:
                    return self._extract_docx(file_path)
                else:
                    print(f"  - DOCX extraction requires python-docx (not installed)")
                    return ""
            elif ext in ['.xlsx', '.xls', '.csv']:
                if PANDAS_AVAILABLE:
                    return self._extract_excel(file_path)
                else:
                    print(f"  - Excel extraction requires pandas (not installed)")
                    return ""
            elif ext in ['.pptx', '.ppt']:
                if PPTX_AVAILABLE:
                    return self._extract_pptx(file_path)
                else:
                    print(f"  - PPTX extraction requires python-pptx (not installed)")
                    return ""
            elif ext in ['.txt', '.md', '.py', '.js', '.json', '.html', '.css', '.java', '.c', '.cpp']:
                return self._extract_plain_text(file_path)
            else:
                print(f"  - Unsupported file type: {ext}")
                return ""
        except Exception as e:
            print(f"  - Error extracting text: {e}")
            return ""

    def _extract_plain_text(self, file_path):
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    def _extract_pdf(self, file_path):
        """Extract PDF with multiple parser fallback for robustness"""
        # Method 1: PyMuPDF (most robust)
        if PYMUPDF_AVAILABLE:
            try:
                doc = fitz.open(str(file_path))
                text = []
                for page in doc:
                    text.append(page.get_text())
                doc.close()
                result = "\n".join(text)
                if result.strip():
                    return result
            except Exception as e:
                print(f"  - PyMuPDF failed: {e}")

        # Method 2: pypdf
        if PYPDF_AVAILABLE:
            try:
                reader = PyPdfReader(str(file_path))
                text = []
                for page in reader.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text.append(extracted)
                result = "\n".join(text)
                if result.strip():
                    return result
            except Exception as e:
                print(f"  - pypdf failed: {e}")

        # Method 3: pdfminer
        if PDFMINER_AVAILABLE:
            try:
                result = pdfminer_extract(str(file_path))
                if result and result.strip():
                    return result
            except Exception as e:
                print(f"  - pdfminer failed: {e}")

        return ""  # All methods failed

    def _extract_docx(self, file_path):
        doc = Document(file_path)
        text = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text.append(cell.text)
        return "\n".join(text)

    def _extract_excel(self, file_path):
        df = pd.read_excel(file_path) if str(file_path).endswith('.xlsx') else pd.read_csv(file_path)
        return df.to_string(index=False, header=True)

    def _extract_pptx(self, file_path):
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

# --- 执行入口 ---
if __name__ == "__main__":
    import sys
    import json

    # 配置
    splitter = DocumentSplitter(
        output_dir="/mnt/agents/chunks",
        chunk_size=32000
    )

    # 支持命令行参数
    if len(sys.argv) > 1:
        file_path = sys.argv[1]
        print(f"Processing file: {file_path}")
        generated_paths = splitter.process_file(file_path)

        # 输出 JSON 格式的结果
        result = {
            "status": "success" if generated_paths else "failed",
            "chunk_files": generated_paths,
            "num_chunks": len(generated_paths)
        }
        print(json.dumps(result, ensure_ascii=False))
    else:
        # 默认示例
        print("Usage: python split_doc.py <file_path>")
        print("Example: python split_doc.py /path/to/document.pdf")